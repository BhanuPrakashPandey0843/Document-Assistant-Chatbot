import os
import time
import base64
import fitz  # PyMuPDF
import concurrent.futures
import docx
import streamlit as st
import spacy
from pptx import Presentation
from langchain.chains import RetrievalQA
from langchain.callbacks.streaming_stdout import StreamingStdOutCallbackHandler
from langchain.callbacks.manager import CallbackManager
from langchain_community.llms import Ollama
from langchain_community.embeddings.ollama import OllamaEmbeddings
from langchain_community.vectorstores import Chroma
from langchain.text_splitter import RecursiveCharacterTextSplitter
from langchain.prompts import PromptTemplate
from langchain.memory import ConversationBufferMemory
import asyncio

# Load SpaCy model
nlp = spacy.load("en_core_web_sm")

# Function to convert image to base64
@st.cache_data
def load_image_to_base64(image_path):
    try:
        with open(image_path, "rb") as image_file:
            return base64.b64encode(image_file.read()).decode()
    except Exception as e:
        st.error(f"Error loading image {image_path}: {e}")
        return ""

# Convert the GIF to base64
bg_image_path = "static/bg.gif"
bg_image_base64 = load_image_to_base64(bg_image_path)

# Page background style
page_bg_img = f"""
<style>
[data-testid="stAppViewContainer"] {{
background-image: url("data:image/gif;base64,{bg_image_base64}");
background-repeat: no-repeat;
background-size: cover;
background-position: center;
background-attachment: fixed;
padding: 0px; /* Optional: to ensure there's no cutting off at the edges */
}}

[data-testid="stVerticalBlock"] {{
color: white;
}}

[data-testid="stHeader"] {{
background-color: rgba(0,0,0,0);
}}

.st-chat {{
    background-color: rgba(255, 0, 0, 1);
    border-radius: 10px;
    padding: 10px;
    margin: 10px 0;
    color: black;
}}

.st-chat .st-chat-message .st-chat-body {{
    font-size: 16px;
    color: black; /* Change the font color to black */
}}

.st-chat .st-chat-avatar {{
    margin-right: 10px;
}}
</style>
"""

# Function to extract text from PDF using PyMuPDF and SpaCy with multithreading
def extract_text_from_pdf(file_path):
    start_time = time.time()
    try:
        st.write(f"Extracting text from PDF: {file_path}")
        document = fitz.open(file_path)
        texts = []

        with concurrent.futures.ThreadPoolExecutor() as executor:
            texts = list(executor.map(lambda page: page.get_text(), document))

        text = " ".join(texts)
        st.write("Text extraction successful.")
        end_time = time.time()
        st.write(f"PDF extraction time: {end_time - start_time:.2f} seconds")
        return text
    except Exception as e:
        st.error(f"Error extracting text from PDF: {e}")
        return ""

# Function to extract text from PowerPoint using SpaCy
def extract_text_from_powerpoint(file_path):
    start_time = time.time()
    try:
        st.write(f"Extracting text from PowerPoint: {file_path}")
        prs = Presentation(file_path)
        texts = []

        def extract_slide_text(slide):
            return " ".join([shape.text_frame.text for shape in slide.shapes if shape.has_text_frame])

        with concurrent.futures.ThreadPoolExecutor() as executor:
            texts = list(executor.map(extract_slide_text, prs.slides))

        text = " ".join(texts)
        st.write("Text extraction successful.")
        end_time = time.time()
        st.write(f"PowerPoint extraction time: {end_time - start_time:.2f} seconds")
        return text
    except Exception as e:
        st.error(f"Error extracting text from PowerPoint: {e}")
        return ""

# Function to extract text from DOC/DOCX files using SpaCy
def extract_text_from_word(file_path):
    start_time = time.time()
    try:
        st.write(f"Extracting text from Word: {file_path}")
        document = docx.Document(file_path)
        text = " ".join([paragraph.text for paragraph in document.paragraphs])
        st.write("Text extraction successful.")
        end_time = time.time()
        st.write(f"Word extraction time: {end_time - start_time:.2f} seconds")
        return text
    except Exception as e:
        st.error(f"Error extracting text from Word: {e}")
        return ""

# Function to extract text from other file types using SpaCy
def extract_text_from_other_file_types(file_path):
    start_time = time.time()
    try:
        st.write(f"Extracting text from other file types: {file_path}")
        with open(file_path, 'r', encoding='utf-8') as file:
            text = file.read()
        st.write("Text extraction successful.")
        end_time = time.time()
        st.write(f"Other file extraction time: {end_time - start_time:.2f} seconds")
        return text
    except Exception as e:
        st.error(f"Error extracting text from file: {e}")
        return ""

# Function to process a document asynchronously
async def process_document(file_path):
    file_extension = os.path.splitext(file_path)[1].lower()
    text = ""

    try:
        if file_extension == '.pdf':
            text = await asyncio.to_thread(extract_text_from_pdf, file_path)
        elif file_extension == '.pptx':
            text = await asyncio.to_thread(extract_text_from_powerpoint, file_path)
        elif file_extension in ['.doc', '.docx']:
            text = await asyncio.to_thread(extract_text_from_word, file_path)
        elif file_extension == '.txt':
            text = await asyncio.to_thread(extract_text_from_other_file_types, file_path)
        else:
            st.error(f"Unsupported file format: {file_extension}")
            return ""

        if text:
            start_time = time.time()
            text_splitter = RecursiveCharacterTextSplitter(
                chunk_size=4000,  # Increased chunk size to reduce the number of chunks
                chunk_overlap=0,  # Removed overlap to minimize redundancy
                length_function=len
            )
            all_chunks = text_splitter.split_text(text)

            st.session_state.vectorstore = Chroma.from_texts(
                texts=all_chunks,
                embedding=OllamaEmbeddings(model="mistral")
            )
            st.session_state.vectorstore.persist()
            end_time = time.time()
            st.write(f"Vectorstore creation time: {end_time - start_time:.2f} seconds")
            st.write("Document processing completed and vectorstore created.")
            return text
        else:
            st.error("Failed to extract text from the document.")
            return ""
    except Exception as e:
        st.error(f"Error processing document: {e}")
        return ""

# Ensure required directories exist
os.makedirs('files', exist_ok=True)
os.makedirs('chroma', exist_ok=True)

# Initialize session state variables
if 'template' not in st.session_state:
    st.session_state.template = """You are a knowledgeable chatbot, here to help with questions of the user. Your tone should be professional and informative.

    Context: {context}
    History: {history}

    User: {question}
    Chatbot:"""

if 'prompt' not in st.session_state:
    st.session_state.prompt = PromptTemplate(
        input_variables=["history", "context", "question"],
        template=st.session_state.template,
    )

if 'memory' not in st.session_state:
    st.session_state.memory = ConversationBufferMemory(
        memory_key="history",
        return_messages=True,
        input_key="question"
    )

if 'vectorstore' not in st.session_state:
    st.session_state.vectorstore = None

if 'llm' not in st.session_state:
    st.session_state.llm = Ollama(model="mistral",
                                  verbose=True,
                                  callback_manager=CallbackManager([StreamingStdOutCallbackHandler()])
                                  )

# Load user and assistant avatars
user_avatar_path = 'user.png'
assistant_avatar_path = 'assistente.png'

# Load avatars in base64
user_avatar_base64 = load_image_to_base64(user_avatar_path)
assistant_avatar_base64 = load_image_to_base64(assistant_avatar_path)

# Initialize chat history
if 'chat_history' not in st.session_state:
    st.session_state.chat_history = []

# Display the page background
st.markdown(page_bg_img, unsafe_allow_html=True)

# Display title and logo
st.title("Document Assistant")
st.sidebar.image('logo.png', use_column_width=True)

# Instructions
st.sidebar.markdown("## Instructions")
st.sidebar.markdown("""
    - **Drag and drop** your document into the designated area or use the upload button below.
    - Once you see a message stating your document has been processed, you can start asking questions in the chat input to interact with the document content.
    ---
""")

# File uploader
st.sidebar.markdown("## Upload File")
uploaded_file = st.sidebar.file_uploader("Upload your file", type=['pdf', 'txt', 'doc', 'docx', 'pptx'])
processar_doc = st.sidebar.button("Process Document")

# Display chat history
for message in st.session_state.chat_history:
    with st.chat_message(message["role"]):
        if message["role"] == "user":
            avatar = f"data:image/png;base64,{user_avatar_base64}"
            st.write             ("User’s question 🗣️")
        elif message["role"] == "assistant":
            st.write("Assistant's response 🤖")
            avatar = f"data:image/png;base64,{assistant_avatar_base64}"
        st.markdown(message["message"])

# Process the uploaded file
if uploaded_file is not None and processar_doc:
    file_path = os.path.join("files", uploaded_file.name)
    with st.spinner("Processing your document..."):
        bytes_data = uploaded_file.read()
        with open(file_path, "wb") as f:
            f.write(bytes_data)

        st.write(f"File {file_path} uploaded successfully.")
        start_time = time.time()
        processed_text = asyncio.run(process_document(file_path))
        end_time = time.time()
        processing_time = end_time - start_time
        st.write(f"Total document processing time: {processing_time:.2f} seconds")

        if processed_text:
            st.success("Your document has been processed.")
            st.balloons()
        else:
            st.error("Document processing failed.")

    if st.session_state.vectorstore:
        st.session_state.retriever = st.session_state.vectorstore.as_retriever()
        if 'qa_chain' not in st.session_state:
            st.session_state.qa_chain = RetrievalQA.from_chain_type(
                llm=st.session_state.llm,
                chain_type='stuff',
                retriever=st.session_state.retriever,
                verbose=True,
                chain_type_kwargs={
                    "verbose": True,
                    "prompt": st.session_state.prompt,
                    "memory": st.session_state.memory,
                }
            )

# Function to handle user input and generate responses
def handle_user_input(user_input):
    user_message = {"role": "user", "message": user_input}
    st.session_state.chat_history.append(user_message)
    with st.chat_message("user", avatar=f"data:image/png;base64,{user_avatar_base64}"):
        st.write("User’s question 🗣️")
        st.markdown(user_input)

    with st.chat_message("assistant", avatar=f"data:image/png;base64,{assistant_avatar_base64}"):
        st.write("Assistant's response 🤖")
        with st.spinner("Assistant is typing..."):
            start_time = time.time()
            response = st.session_state.qa_chain(user_input)
            end_time = time.time()
            response_time = end_time - start_time
            st.write(f"Response generation time: {response_time:.2f} seconds")

        message_placeholder = st.empty()
        full_response = ""
        for chunk in response['result'].split():
            full_response += chunk + " "
            time.sleep(0.05)
            message_placeholder.markdown(full_response + "▌")
        message_placeholder.markdown(full_response)

    chatbot_message = {"role": "assistant", "message": response['result']}
    st.session_state.chat_history.append(chatbot_message)

# Handle user input
if st.session_state.vectorstore:
    user_input = st.chat_input("You:", key="user_input")
    if user_input:
        handle_user_input(user_input)
else:
    st.warning("Please upload a file and click the process button.", icon="⚠️")

